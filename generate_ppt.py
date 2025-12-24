import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Define colors based on design prompt
COLOR_BG = RGBColor(255, 255, 255)       # White
COLOR_ACCENT_1 = RGBColor(2, 132, 199)   # Sky 600
COLOR_ACCENT_2 = RGBColor(79, 70, 229)   # Indigo 600
COLOR_TEXT_TITLE = RGBColor(15, 23, 42)  # Slate 900
COLOR_TEXT_BODY = RGBColor(71, 85, 105)  # Slate 600
COLOR_WHITE = RGBColor(255, 255, 255)

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    create_title_slide(prs, 
        title="냉난방기 스마트 정비시스템\n프로세스 설명서", 
        subtitle="현업 담당자용 가이드 | Ver 1.0 (2025.12.13)",
        footer="작성: DMES팀"
    )

    # 2. Introduction
    create_content_slide(prs,
        title="1. 시스템 소개",
        content_title="시스템 개요",
        content_body="공조설비의 고장신고부터 정비완료까지 전 과정을 디지털화하여 효율적인 정비를 지원합니다.\n\n" +
                     "- QR코드 스캔으로 빠른 신고\n" +
                     "- 실시간 진행상황 확인 (신고~승인~완료)\n" +
                     "- 모바일/웹 완벽 연동"
    )

    # 3. Key Features (Grid Layout)
    create_grid_slide(prs, "1.2 주요 특징", [
        ("QR코드 기반", "설비 QR 스캔으로 장비 정보 자동 로드.\n장비코드/위치 입력 불필요."),
        ("실시간 알림", "주요 단계별 담당자 자동 알림.\n(신고, 승인, 시작, 완료 등)"),
        ("모바일/웹 연동", "현장: 모바일 앱\n사무실: 웹 시스템\n언제 어디서나 업무 처리 가능."),
        ("백업 기능", "QR 미인식/네트워크 오류 등\n예외 상황 시 관리자 웹 대행 처리.")
    ])

    # 4. Process Overview
    create_content_slide(prs,
        title="2. 전체 프로세스 요약",
        content_title="4단계 정비 프로세스",
        content_body="1. [Step 1] 고장신고 & 승인 (현장/승인권자)\n" +
                     "2. [Step 2] 방문예정일 등록 (정비접수자)\n" +
                     "3. [Step 3] 작업 시작 & 완료 (정비담당자)\n" +
                     "4. [Step 4] 작업 종료 & 평가 (관리자)"
    )

    # 5. Step 1 Detail
    create_content_slide(prs,
        title="3. Step 1: 고장신고 & 승인",
        content_title="담당자: 신고담당자 / 승인권자",
        content_body="[고장신고]\n" +
                     "1. 앱 실행 > '고장' 또는 'Overhaul' 선택\n" +
                     "2. QR코드 1차 스캔 (설비 정보 자동 로드)\n" +
                     "3. 증상 입력 및 전송\n\n" +
                     "[승인처리]\n" +
                     "- 알림 수신 후 승인 또는 반려 처리\n" +
                     "- 승인 시 방문요청일 지정"
    )

    # 6. Step 2 Detail
    create_content_slide(prs,
        title="4. Step 2: 방문예정일 등록",
        content_title="담당자: 정비접수자 (협력사)",
        content_body="1. 접수 내역 확인 및 담당자 지정\n" +
                     "2. 방문 예정일 입력\n" +
                     "3. 신고자에게 일정 알림 자동 발송\n\n" +
                     "* 일정 변경 시 사유 등록 후 수정 (재알림 발송)"
    )

    # 7. Step 3 Detail
    create_content_slide(prs,
        title="5. Step 3: 작업 시작 & 완료",
        content_title="담당자: 정비담당자 (협력사)",
        content_body="[작업 시작]\n" +
                     "- 현장 도착 후 QR코드 2차 스캔\n" +
                     "- '작업시작' 등록 (알림 발송)\n\n" +
                     "[작업 완료]\n" +
                     "- 정비 후 QR코드 3차 스캔\n" +
                     "- 조치내용, 부품, 시간 입력 후 '작업완료' 등록"
    )

    # 8. Step 4 Detail
    create_content_slide(prs,
        title="6. Step 4: 작업 종료 & 평가",
        content_title="담당자: 관리자 (설비팀)",
        content_body="1. 작업 완료 알림 수신\n" +
                     "2. 완료 대상 확인 및 만족도 평가 (선택사항)\n" +
                     "3. 정비 완료 처리 (이력 자동 저장)\n\n" +
                     "* 필요 시 취소 사유 입력 후 강제 종료 가능"
    )

    # 9. Roles & Responsibilities
    create_grid_slide(prs, "7. 사용자별 역할", [
        ("신고담당자", "고장 발견/신고\n진행현황 조회"),
        ("승인권자", "정비 타당성 검토\n승인/반려"),
        ("정비담당자", "실제 수리 수행\n시작/완료 QR처리"),
        ("관리자", "운영 총괄\n백업/평가/통계")
    ])

    # 10. QR Scan Points
    create_content_slide(prs,
        title="8. QR코드 활용 (총 3회)",
        content_title="스캔 시점 및 목적",
        content_body="1. [1차] 고장 신고 시\n   - 장비정보 로드, 신고 접수\n\n" +
                     "2. [2차] 작업 시작 시\n   - 현장 도착 확인, 작업 시작 알림\n\n" +
                     "3. [3차] 작업 완료 시\n   - 작업 종료 확인, 결과 입력 활성화"
    )

    # 11. Backup Functions
    create_content_slide(prs,
        title="10. 백업(예외처리) 기능",
        content_title="관리자 웹을 통한 대행 처리",
        content_body="- 요청 백업: QR/앱 불가 시 대리 신고\n" +
                     "- 취소 백업: 잘못된 접수 강제 취소\n" +
                     "- 승인 백업: 승인권자 부재 시 대리 승인\n" +
                     "- 종료 백업: 정상 종료 누락 건 일괄 처리\n" +
                     "- 예정일 변경: 시스템 오류 시 일정 수정"
    )

    # Save
    filename = '냉난방기_스마트정비_시스템.pptx'
    prs.save(filename)
    print(f"Presentation saved to {filename}")

def create_title_slide(prs, title, subtitle, footer):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Customize background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Style Title
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_1

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle + "\n" + footer
    
    # Style Subtitle
    for p in subtitle_shape.text_frame.paragraphs:
        p.font.name = 'Malgun Gothic'
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXT_BODY

def create_content_slide(prs, title, content_title, content_body):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_TITLE
    p.alignment = PP_ALIGN.LEFT

    # Content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content_title
    
    p = tf.paragraphs[0]
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_1
    
    p_body = tf.add_paragraph()
    p_body.text = content_body
    p_body.font.name = 'Malgun Gothic'
    p_body.font.size = Pt(20)
    p_body.font.color.rgb = COLOR_TEXT_BODY
    p_body.line_spacing = 1.5

def create_grid_slide(prs, title, items):
    """
    items: list of tuples (header, body)
    Creates a 2x2 or 1x4 grid effect using shapes
    """
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG

    # Title manually added
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_TITLE

    # Grid config
    start_x = Inches(0.5)
    start_y = Inches(1.5)
    width = Inches(4.25)
    height = Inches(2.2)
    margin = Inches(0.2)

    for i, (header, body) in enumerate(items):
        row = i // 2
        col = i % 2
        x = start_x + (col * (width + margin))
        y = start_y + (row * (height + margin))
        
        # Draw box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_WHITE
        shape.line.color.rgb = COLOR_ACCENT_2
        shape.line.width = Pt(1.5)
        
        # Add shadow text frame specifically? No, just use shape text
        tf = shape.text_frame
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = header
        p.font.name = 'Malgun Gothic'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT_2
        
        p_body = tf.add_paragraph()
        p_body.text = body
        p_body.font.name = 'Malgun Gothic'
        p_body.font.size = Pt(16)
        p_body.font.color.rgb = COLOR_TEXT_BODY
        p_body.space_before = Pt(10)

if __name__ == "__main__":
    create_presentation()
