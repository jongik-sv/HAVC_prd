#!/usr/bin/env python3
"""
공조설비 정보관리 시스템 구축 - 사전 설명회 PPTX 생성기

PPT기본양식.pptx 템플릿을 기반으로 presentation_content.json 내용을 채워 넣어
최종 PPTX 파일을 생성합니다.

레이아웃 매핑:
- Layout 0: White_Big K 버전 (표지) - layout_id: 1
- Layout 1: 간지 1 (목차) - layout_id: 2
- Layout 2: 내지 (Action title 사용) - layout_id: 3
- Layout 3: 내지 (Action title, Body삭제) - layout_id: 4
- Layout 4: 내지 (Action Title 미사용) - layout_id: 5
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# ==================== 경로 설정 ====================
BASE_DIR = Path(__file__).parent
TEMPLATE_PATH = BASE_DIR / "PPT기본양식.pptx"
CONTENT_PATH = BASE_DIR / "presentation_content.json"
OUTPUT_PATH = BASE_DIR / "공조설비_정보관리시스템_사전설명회.pptx"

# ==================== 색상 정의 ====================
COLORS = {
    "navy": RGBColor(0x00, 0x24, 0x52),
    "red": RGBColor(0xC5, 0x1F, 0x2A),
    "gray": RGBColor(0x66, 0x66, 0x66),
    "light_gray": RGBColor(0x99, 0x99, 0x99),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "green": RGBColor(0x2E, 0x7D, 0x32),
    "orange": RGBColor(0xF5, 0x7C, 0x00),
    "blue": RGBColor(0x19, 0x76, 0xD2),
}

# ==================== 레이아웃 인덱스 매핑 ====================
# PPT기본양식.pptx의 실제 레이아웃 순서
LAYOUT_MAP = {
    1: 0,  # 표지 (White_Big K)
    2: 1,  # 목차 (간지 1)
    3: 2,  # 내지 (Action title 사용)
    4: 3,  # 내지 (Action title, Body삭제)
    5: 4,  # 내지 (Action Title 미사용)
}


def load_content():
    """JSON 콘텐츠 파일 로드"""
    with open(CONTENT_PATH, "r", encoding="utf-8") as f:
        return json.load(f)


def get_placeholder_by_idx(slide, idx):
    """플레이스홀더 인덱스로 shape 찾기"""
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.idx == idx:
                return shape
    return None


def set_text_frame_content(text_frame, text, font_size=None, bold=False, color=None):
    """텍스트 프레임에 내용 설정"""
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color


def add_title_slide(prs, slide_data):
    """슬라이드 1: 표지 추가"""
    layout = prs.slide_layouts[LAYOUT_MAP[1]]
    slide = prs.slides.add_slide(layout)

    placeholders = slide_data.get("placeholders", {})

    # 제목 설정 (idx=15 또는 title placeholder)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            idx = shape.placeholder_format.idx

            # 제목 플레이스홀더
            if idx == 15 or ph_type == 1:  # title type
                if shape.has_text_frame:
                    set_text_frame_content(
                        shape.text_frame,
                        placeholders.get("title", ""),
                        font_size=32,
                        bold=True,
                        color=COLORS["navy"]
                    )

            # 부제목 플레이스홀더 (idx=13 또는 27)
            elif idx in [13, 27]:
                if shape.has_text_frame:
                    set_text_frame_content(
                        shape.text_frame,
                        placeholders.get("subtitle", ""),
                        font_size=14,
                        color=COLORS["gray"]
                    )

    return slide


def add_toc_slide(prs, slide_data):
    """슬라이드 2: 목차 추가"""
    layout = prs.slide_layouts[LAYOUT_MAP[2]]
    slide = prs.slides.add_slide(layout)

    placeholders = slide_data.get("placeholders", {})
    toc_items = placeholders.get("toc_items", [])

    # 목차 플레이스홀더들 찾기 (idx=17: 번호, idx=13: 제목, idx=14: 페이지)
    number_ph = get_placeholder_by_idx(slide, 17)
    title_ph = get_placeholder_by_idx(slide, 13)
    page_ph = get_placeholder_by_idx(slide, 14)

    # 번호 열
    if number_ph and number_ph.has_text_frame:
        tf = number_ph.text_frame
        tf.clear()
        for i, item in enumerate(toc_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item.get("number", "")
            p.font.size = Pt(16)
            p.font.bold = True
            p.alignment = PP_ALIGN.RIGHT

    # 제목 열
    if title_ph and title_ph.has_text_frame:
        tf = title_ph.text_frame
        tf.clear()
        for i, item in enumerate(toc_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item.get("title", "")
            p.font.size = Pt(16)

    # 페이지 열
    if page_ph and page_ph.has_text_frame:
        tf = page_ph.text_frame
        tf.clear()
        for i, item in enumerate(toc_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item.get("pages", "")
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS["light_gray"]

    return slide


def add_content_slide_layout4(prs, slide_data):
    """Layout 4: 내지 (Action title, Body삭제) - 자유 콘텐츠용"""
    layout = prs.slide_layouts[LAYOUT_MAP[4]]
    slide = prs.slides.add_slide(layout)

    placeholders = slide_data.get("placeholders", {})

    # Main Title (idx=19)
    main_title_ph = get_placeholder_by_idx(slide, 19)
    if main_title_ph and main_title_ph.has_text_frame:
        set_text_frame_content(
            main_title_ph.text_frame,
            placeholders.get("main_title", ""),
            font_size=19,
            bold=False,
            color=COLORS["navy"]
        )

    # Action Title (title placeholder)
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.type == 1:  # title
                if shape.has_text_frame:
                    set_text_frame_content(
                        shape.text_frame,
                        placeholders.get("action_title", ""),
                        font_size=17,
                        color=COLORS["gray"]
                    )
                break

    # custom_elements 처리는 별도 함수에서
    custom_elements = slide_data.get("custom_elements", [])
    add_custom_elements(slide, custom_elements)

    return slide


def add_content_slide_layout5(prs, slide_data):
    """Layout 5: 내지 (Action Title 미사용) - 넓은 본문용"""
    layout = prs.slide_layouts[LAYOUT_MAP[5]]
    slide = prs.slides.add_slide(layout)

    placeholders = slide_data.get("placeholders", {})

    # Main Title (idx=19)
    main_title_ph = get_placeholder_by_idx(slide, 19)
    if main_title_ph and main_title_ph.has_text_frame:
        set_text_frame_content(
            main_title_ph.text_frame,
            placeholders.get("main_title", ""),
            font_size=19,
            bold=False,
            color=COLORS["navy"]
        )

    # Body (idx=18)
    body_ph = get_placeholder_by_idx(slide, 18)
    if body_ph and body_ph.has_text_frame:
        body_items = placeholders.get("body", [])
        tf = body_ph.text_frame
        tf.clear()

        for i, item in enumerate(body_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item.get("text", "")
            p.font.size = Pt(16)
            p.level = item.get("level", 1) - 1

    return slide


def add_custom_elements(slide, elements):
    """커스텀 요소들 추가 (표, 차트, 다이어그램 등)"""
    # 콘텐츠 영역 시작 위치
    content_top = Emu(1431130)  # 약 4cm
    content_left = Emu(270064)  # 약 0.75cm
    content_width = Emu(9360550)  # 약 26cm

    for element in elements:
        elem_type = element.get("type", "")
        data = element.get("data", {})

        if elem_type == "table":
            add_table_element(slide, data, content_left, content_top, content_width)
        elif elem_type == "icon_box_grid":
            add_icon_box_grid(slide, data, content_left, content_top, content_width)
        elif elem_type == "pain_point_cards":
            add_pain_point_cards(slide, data, content_left, content_top, content_width)
        elif elem_type == "process_flow":
            add_process_flow(slide, data, content_left, content_top, content_width)
        elif elem_type == "comparison_chart":
            add_comparison_chart(slide, data, content_left, content_top, content_width)
        elif elem_type == "timeline":
            add_timeline(slide, data, content_left, content_top, content_width)
        elif elem_type == "screen_gallery":
            add_screen_gallery(slide, data, content_left, content_top, content_width)
        # architecture_diagram은 복잡하므로 placeholder로 처리


def add_table_element(slide, data, left, top, width):
    """테이블 요소 추가"""
    headers = data.get("headers", [])
    rows = data.get("rows", [])

    if not headers or not rows:
        return

    cols = len(headers)
    row_count = len(rows) + 1  # 헤더 포함

    # 테이블 크기 계산
    table_height = Emu(row_count * 500000)  # 행당 약 1.4cm

    table = slide.shapes.add_table(row_count, cols, left, top, width, table_height).table

    # 헤더 행 스타일
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["navy"]

        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = COLORS["white"]
        para.alignment = PP_ALIGN.LEFT

    # 데이터 행
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)

            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(13)
            para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            # 심각도 열 특별 처리
            if col_idx == len(row_data) - 1 and cell_text in ["높음", "중간", "낮음"]:
                if cell_text == "높음":
                    para.font.bold = True
                    para.font.color.rgb = COLORS["red"]
                elif cell_text == "중간":
                    para.font.bold = True
                    para.font.color.rgb = COLORS["orange"]


def add_icon_box_grid(slide, data, left, top, width):
    """아이콘 박스 그리드 추가"""
    items = data.get("items", [])
    columns = data.get("columns", 4)

    if not items:
        return

    box_width = Emu(2000000)  # 약 5.5cm
    box_height = Emu(1800000)  # 약 5cm
    gap = Emu(200000)  # 약 0.5cm

    accent_colors = [COLORS["navy"], COLORS["red"], COLORS["gray"], COLORS["orange"]]

    for i, item in enumerate(items):
        col = i % columns
        row = i // columns

        x = left + (box_width + gap) * col
        y = top + (box_height + gap) * row

        # 카드 배경 (사각형)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["white"]
        shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        # 좌측 컬러 바
        color_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y, Emu(50000), box_height
        )
        color_bar.fill.solid()
        color_bar.fill.fore_color.rgb = accent_colors[i % len(accent_colors)]
        color_bar.line.fill.background()

        # 타이틀 텍스트박스
        title_box = slide.shapes.add_textbox(
            x + Emu(150000), y + Emu(800000),
            box_width - Emu(200000), Emu(300000)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.get("title", "")
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["navy"]

        # 설명 텍스트박스
        desc_box = slide.shapes.add_textbox(
            x + Emu(150000), y + Emu(1100000),
            box_width - Emu(200000), Emu(600000)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.get("desc", "").replace("\\n", "\n")
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["gray"]


def add_pain_point_cards(slide, data, left, top, width):
    """Pain Point 카드 추가"""
    items = data.get("items", [])
    columns = data.get("columns", 4)

    if not items:
        return

    box_width = Emu(2000000)
    box_height = Emu(2200000)
    gap = Emu(150000)

    role_colors = [
        COLORS["blue"],
        COLORS["green"],
        COLORS["orange"],
        RGBColor(0x7B, 0x1F, 0xA2)  # purple
    ]

    for i, item in enumerate(items):
        col = i % columns
        x = left + (box_width + gap) * col
        y = top

        # 카드 배경
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
        shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        # 역할 라벨
        role_box = slide.shapes.add_textbox(
            x + Emu(100000), y + Emu(700000),
            box_width - Emu(200000), Emu(300000)
        )
        tf = role_box.text_frame
        p = tf.paragraphs[0]
        p.text = item.get("role", "")
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = role_colors[i % len(role_colors)]
        p.alignment = PP_ALIGN.CENTER

        # Pain 박스 배경
        pain_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Emu(100000), y + Emu(1100000),
            box_width - Emu(200000), Emu(900000)
        )
        pain_bg.fill.solid()
        pain_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
        pain_bg.line.fill.background()

        # Pain 텍스트
        pain_box = slide.shapes.add_textbox(
            x + Emu(200000), y + Emu(1200000),
            box_width - Emu(400000), Emu(700000)
        )
        tf = pain_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.get("pain", "").replace("\\n", "\n")
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["red"]


def add_process_flow(slide, data, left, top, width):
    """프로세스 플로우 추가"""
    steps = data.get("steps", [])

    if not steps:
        return

    circle_size = Emu(900000)  # 약 2.5cm
    gap = Emu(300000)
    arrow_width = Emu(400000)

    total_width = len(steps) * circle_size + (len(steps) - 1) * (gap + arrow_width)
    start_x = left + (width - total_width) // 2

    step_colors = {
        "navy": COLORS["navy"],
        "green": COLORS["green"],
        "orange": COLORS["orange"],
    }

    for i, step in enumerate(steps):
        x = start_x + i * (circle_size + gap + arrow_width)
        y = top + Emu(300000)

        # 원형 노드
        color_key = "navy"
        if "내부" in step.get("actor", ""):
            color_key = "navy"
        elif "승인" in step.get("name", ""):
            color_key = "green"
        else:
            color_key = "orange"

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, y, circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = step_colors[color_key]
        circle.line.fill.background()

        # 상태 코드
        code_box = slide.shapes.add_textbox(
            x, y + Emu(300000),
            circle_size, Emu(400000)
        )
        tf = code_box.text_frame
        p = tf.paragraphs[0]
        p.text = step.get("code", "")
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER

        # 단계명
        name_box = slide.shapes.add_textbox(
            x - Emu(200000), y + circle_size + Emu(100000),
            circle_size + Emu(400000), Emu(300000)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = step.get("name", "")
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = step_colors[color_key]
        p.alignment = PP_ALIGN.CENTER

        # 담당자
        actor_box = slide.shapes.add_textbox(
            x - Emu(200000), y + circle_size + Emu(350000),
            circle_size + Emu(400000), Emu(250000)
        )
        tf = actor_box.text_frame
        p = tf.paragraphs[0]
        p.text = step.get("actor", "")
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS["light_gray"]
        p.alignment = PP_ALIGN.CENTER

        # 화살표 (마지막 제외)
        if i < len(steps) - 1:
            arrow_x = x + circle_size + Emu(50000)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_x, y + circle_size // 2 - Emu(100000),
                arrow_width, Emu(200000)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            arrow.line.fill.background()


def add_comparison_chart(slide, data, left, top, width):
    """비교 차트 추가"""
    items = data.get("items", [])

    if not items:
        return

    row_height = Emu(650000)
    label_width = Emu(2000000)
    bar_area_width = Emu(5500000)
    change_width = Emu(1200000)

    for i, item in enumerate(items):
        y = top + row_height * i

        # 라벨
        label_box = slide.shapes.add_textbox(
            left, y + Emu(150000),
            label_width, Emu(400000)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = item.get("label", "")
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # As-Is 바
        as_is_value = item.get("as_is", 0)
        max_value = max(item.get("as_is", 0), item.get("to_be", 0), 1)
        as_is_width = int(bar_area_width * (as_is_value / 100)) if item.get("unit") == "%" else int(bar_area_width * (as_is_value / max_value / 1.2))

        bar_x = left + label_width + Emu(200000)
        as_is_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            bar_x, y + Emu(50000),
            max(as_is_width, Emu(100000)), Emu(200000)
        )
        as_is_bar.fill.solid()
        as_is_bar.fill.fore_color.rgb = RGBColor(0xB0, 0xBE, 0xC5)
        as_is_bar.line.fill.background()

        # To-Be 바
        to_be_value = item.get("to_be", 0)
        to_be_width = int(bar_area_width * (to_be_value / 100)) if item.get("unit") == "%" else int(bar_area_width * (to_be_value / max_value / 1.2))

        to_be_color = COLORS["green"] if "+" in item.get("change", "") else COLORS["navy"]
        to_be_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            bar_x, y + Emu(300000),
            max(to_be_width, Emu(100000)), Emu(200000)
        )
        to_be_bar.fill.solid()
        to_be_bar.fill.fore_color.rgb = to_be_color
        to_be_bar.line.fill.background()

        # 변화율 뱃지
        change_box = slide.shapes.add_textbox(
            left + label_width + bar_area_width + Emu(400000), y + Emu(150000),
            change_width, Emu(300000)
        )
        tf = change_box.text_frame
        p = tf.paragraphs[0]
        change_text = item.get("change", "")
        p.text = f"▼{change_text[1:]}" if change_text.startswith("-") else f"▲{change_text[1:]}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS["green"] if "+" in change_text else COLORS["blue"]
        p.alignment = PP_ALIGN.RIGHT


def add_timeline(slide, data, left, top, width):
    """타임라인 추가 - 슬라이드에 맞게 축소"""
    phases = data.get("phases", [])

    if not phases:
        return

    # 타임라인 전체 폭을 슬라이드 70%로 제한
    timeline_width = Emu(7000000)  # 약 19.4cm
    timeline_left = left + (width - timeline_width) // 2

    # 주차 헤더
    week_width = timeline_width // 12
    header_height = Emu(300000)

    for week in range(12):
        x = timeline_left + week_width * week

        header_box = slide.shapes.add_textbox(
            x, top, week_width, header_height
        )
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{week + 1}주"
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS["light_gray"]
        p.alignment = PP_ALIGN.CENTER

    # Phase 바
    bar_height = Emu(350000)
    phase_colors = [COLORS["navy"], COLORS["red"]]

    for i, phase in enumerate(phases):
        y = top + header_height + Emu(150000) + (bar_height + Emu(250000)) * i

        # Phase 라벨
        label_width = Emu(900000)
        label_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            timeline_left - label_width - Emu(100000), y,
            label_width, bar_height
        )
        label_box.fill.solid()
        label_box.fill.fore_color.rgb = phase_colors[i]
        label_box.line.fill.background()

        label_text = slide.shapes.add_textbox(
            timeline_left - label_width - Emu(80000), y + Emu(80000),
            label_width - Emu(40000), Emu(200000)
        )
        tf = label_text.text_frame
        p = tf.paragraphs[0]
        p.text = phase.get("name", "").split(" - ")[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER

        # Phase 바
        if i == 0:
            bar_start = timeline_left
            bar_width = week_width * 8
        else:
            bar_start = timeline_left + week_width * 8
            bar_width = week_width * 4

        phase_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            bar_start, y,
            bar_width, bar_height
        )
        phase_bar.fill.solid()
        phase_bar.fill.fore_color.rgb = phase_colors[i]
        phase_bar.line.fill.background()


def add_screen_gallery(slide, data, left, top, width):
    """스크린샷 갤러리 추가 (모바일/웹 화면 이미지) - 원본 비율 유지"""
    from pathlib import Path
    from PIL import Image

    screens = data.get("screens", [])
    layout_type = data.get("layout", "horizontal_3")

    if not screens:
        return

    num_screens = len(screens)
    is_wide = "wide" in layout_type.lower()

    # 최대 높이 제한 (슬라이드 영역 고려)
    if is_wide:
        max_height = Emu(2600000)  # 웹: 약 7.2cm
    else:
        max_height = Emu(3000000)  # 모바일: 약 8.3cm

    gap = Emu(200000)  # 이미지 간격

    # 먼저 각 이미지의 실제 크기를 계산
    image_infos = []
    for screen in screens:
        image_path = screen.get("image_path", "")
        if image_path and Path(image_path).exists():
            try:
                with Image.open(image_path) as img:
                    orig_width, orig_height = img.size
                    # 높이 기준으로 비율 계산
                    scale = int(max_height) / orig_height
                    new_width = Emu(int(orig_width * scale))
                    new_height = max_height
                    image_infos.append({
                        "path": image_path,
                        "width": new_width,
                        "height": new_height,
                        "label": screen.get("label", ""),
                        "description": screen.get("description", "").replace("\\n", "\n")
                    })
            except Exception as e:
                print(f"    ⚠️  이미지 정보 로드 실패: {image_path} - {e}")
                # 기본 크기 사용
                default_width = Emu(2000000) if is_wide else Emu(1600000)
                image_infos.append({
                    "path": image_path,
                    "width": default_width,
                    "height": max_height,
                    "label": screen.get("label", ""),
                    "description": screen.get("description", "").replace("\\n", "\n")
                })
        else:
            default_width = Emu(2800000) if is_wide else Emu(1600000)
            image_infos.append({
                "path": None,
                "width": default_width,
                "height": max_height,
                "label": screen.get("label", ""),
                "description": screen.get("description", "").replace("\\n", "\n")
            })

    # 전체 너비 계산 및 시작 위치 중앙 정렬
    total_width = sum(info["width"] for info in image_infos) + (num_screens - 1) * gap
    start_x = left + (width - total_width) // 2

    current_x = start_x
    for info in image_infos:
        x = current_x
        y = top
        img_width = info["width"]
        img_height = info["height"]

        if info["path"] and Path(info["path"]).exists():
            try:
                pic = slide.shapes.add_picture(
                    info["path"],
                    x, y,
                    width=img_width,
                    height=img_height
                )
                pic.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                pic.line.width = Pt(1)
            except Exception as e:
                print(f"    ⚠️  이미지 로드 실패: {info['path']} - {e}")
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x, y, img_width, img_height
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
                placeholder.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        else:
            print(f"    ⚠️  이미지 파일 없음: {info['path']}")
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, img_width, img_height
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            placeholder.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        # 라벨 (이미지 아래)
        label_box = slide.shapes.add_textbox(
            x, y + img_height + Emu(80000),
            img_width, Emu(250000)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = info["label"]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS["navy"]
        p.alignment = PP_ALIGN.CENTER

        # 설명 텍스트 (라벨 아래)
        if info["description"]:
            desc_box = slide.shapes.add_textbox(
                x, y + img_height + Emu(300000),
                img_width, Emu(400000)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = info["description"]
            p.font.size = Pt(9)
            p.font.color.rgb = COLORS["gray"]
            p.alignment = PP_ALIGN.CENTER

        current_x += img_width + gap


def generate_presentation():
    """메인 프레젠테이션 생성 함수"""
    print("=" * 50)
    print("PPTX 생성 시작")
    print("=" * 50)

    # 템플릿 로드
    print(f"템플릿 로드: {TEMPLATE_PATH}")
    prs = Presentation(str(TEMPLATE_PATH))

    # 기존 슬라이드 삭제 (템플릿 샘플 슬라이드)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # 콘텐츠 로드
    print(f"콘텐츠 로드: {CONTENT_PATH}")
    content = load_content()

    slides_data = content.get("presentation", {}).get("slides", [])

    # 슬라이드 생성
    for slide_data in slides_data:
        slide_num = slide_data.get("slide_number", 0)
        layout_id = slide_data.get("layout_id", 4)

        print(f"  슬라이드 {slide_num} 생성 (layout_id: {layout_id})")

        if layout_id == 1:
            add_title_slide(prs, slide_data)
        elif layout_id == 2:
            add_toc_slide(prs, slide_data)
        elif layout_id == 4:
            add_content_slide_layout4(prs, slide_data)
        elif layout_id == 5:
            add_content_slide_layout5(prs, slide_data)
        else:
            # 기본적으로 layout 4 사용
            add_content_slide_layout4(prs, slide_data)

    # 저장
    print(f"\n저장: {OUTPUT_PATH}")
    prs.save(str(OUTPUT_PATH))
    print("=" * 50)
    print("PPTX 생성 완료!")
    print("=" * 50)


if __name__ == "__main__":
    generate_presentation()
