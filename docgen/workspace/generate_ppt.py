#!/usr/bin/env python3
"""
Generate PPT from presentation_content.json using template
"""
import json
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def load_content(json_path):
    """Load presentation content from JSON"""
    with open(json_path, 'r', encoding='utf-8') as f:
        return json.load(f)

def find_placeholder_by_type(slide, ph_type):
    """Find placeholder by type name"""
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.type.name == ph_type:
                return shape
    return None

def find_shape_by_name(slide, name_contains):
    """Find shape by partial name match"""
    for shape in slide.shapes:
        if name_contains.lower() in shape.name.lower():
            return shape
    return None

def clear_text_frame(shape):
    """Clear all text from a text frame"""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = ""

def set_text_with_style(shape, text, font_size=None, bold=False, color=None):
    """Set text on a shape with styling"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Clear existing
    for para in tf.paragraphs:
        para.clear()
    # Set new text
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor.from_string(color)

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color="000000"):
    """Add a text box to a slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return txBox

def populate_title_slide(slide, content):
    """Populate title slide (slide 1)"""
    placeholders = content.get('placeholders', {})
    title = placeholders.get('title', '')
    subtitle = placeholders.get('subtitle', '')

    # Find title and subtitle placeholders
    title_shape = find_placeholder_by_type(slide, 'TITLE') or find_placeholder_by_type(slide, 'CENTER_TITLE')
    subtitle_shape = find_placeholder_by_type(slide, 'SUBTITLE')

    if title_shape and title:
        if title_shape.has_text_frame:
            title_shape.text_frame.paragraphs[0].text = title

    if subtitle_shape and subtitle:
        if subtitle_shape.has_text_frame:
            subtitle_shape.text_frame.paragraphs[0].text = subtitle

def populate_toc_slide(slide, content):
    """Populate table of contents slide (slide 2)"""
    placeholders = content.get('placeholders', {})
    toc_items = placeholders.get('toc_items', [])

    # Find body placeholders
    body_shapes = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.type.name == 'BODY':
                body_shapes.append(shape)

    # Sort by position (left to right)
    body_shapes.sort(key=lambda s: s.left)

    if len(body_shapes) >= 3 and toc_items:
        # shape-0: numbers, shape-1: titles, shape-2: pages
        numbers_shape = body_shapes[0]
        titles_shape = body_shapes[1]
        pages_shape = body_shapes[2]

        # Clear and populate
        if numbers_shape.has_text_frame:
            tf = numbers_shape.text_frame
            for i, item in enumerate(toc_items[:9]):  # Max 9 items
                if i < len(tf.paragraphs):
                    tf.paragraphs[i].text = item.get('number', '')
                else:
                    p = tf.add_paragraph()
                    p.text = item.get('number', '')

        if titles_shape.has_text_frame:
            tf = titles_shape.text_frame
            for i, item in enumerate(toc_items[:9]):
                if i < len(tf.paragraphs):
                    tf.paragraphs[i].text = item.get('title', '')
                else:
                    p = tf.add_paragraph()
                    p.text = item.get('title', '')

        if pages_shape.has_text_frame:
            tf = pages_shape.text_frame
            for i, item in enumerate(toc_items[:9]):
                if i < len(tf.paragraphs):
                    tf.paragraphs[i].text = item.get('pages', '')
                else:
                    p = tf.add_paragraph()
                    p.text = item.get('pages', '')

def populate_content_slide(slide, content, slide_idx):
    """Populate content slide with main_title, action_title, and body if available"""
    placeholders = content.get('placeholders', {})
    main_title = placeholders.get('main_title', '')
    action_title = placeholders.get('action_title', '')
    body = placeholders.get('body', [])

    # Find title placeholder
    title_shape = find_placeholder_by_type(slide, 'TITLE')
    if title_shape and title_shape.has_text_frame:
        tf = title_shape.text_frame
        # Clear existing paragraphs
        for para in tf.paragraphs:
            para.clear()

        # Add main title if exists
        if main_title:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = main_title
            run.font.bold = True
            run.font.size = Pt(20)

        # Add action title on new line if exists
        if action_title:
            if main_title:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            run = p.add_run()
            run.text = action_title
            run.font.size = Pt(14)

    # Handle body content if present (layout_id 3 slides)
    if body:
        body_shape = find_placeholder_by_type(slide, 'BODY')
        if body_shape and body_shape.has_text_frame:
            tf = body_shape.text_frame
            # Clear existing
            for para in tf.paragraphs:
                para.clear()

            # Add body items
            for i, item in enumerate(body):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                level = item.get('level', 1)
                text = item.get('text', '')

                p.text = text
                p.level = level - 1  # PowerPoint levels are 0-indexed

                if level == 1:
                    p.font.bold = True
                    p.font.size = Pt(14)
                else:
                    p.font.size = Pt(12)

def update_footer(slide, footer_text):
    """Update footer on slide"""
    footer_shape = find_placeholder_by_type(slide, 'FOOTER')
    if footer_shape and footer_shape.has_text_frame:
        footer_shape.text_frame.paragraphs[0].text = footer_text

def main():
    # Paths
    template_path = r"c:\Project\HAVC_prd\docgen\workspace\working.pptx"
    content_path = r"c:\Project\HAVC_prd\docgen\presentation_content.json"
    output_path = r"c:\Project\HAVC_prd\docgen\presentation_output.pptx"

    # Load content
    data = load_content(content_path)
    presentation_data = data.get('presentation', {})
    slides_content = presentation_data.get('slides', [])
    author = presentation_data.get('author', '작성자_동국시스템즈')

    # Open template
    prs = Presentation(template_path)

    # Process each slide
    for i, slide in enumerate(prs.slides):
        if i >= len(slides_content):
            break

        content = slides_content[i]
        layout_id = content.get('layout_id', 4)

        print(f"Processing slide {i+1} (layout_id: {layout_id})")

        if i == 0:  # Title slide
            populate_title_slide(slide, content)
        elif i == 1:  # TOC slide
            populate_toc_slide(slide, content)
        else:  # Content slides
            populate_content_slide(slide, content, i)

        # Update footer with author
        update_footer(slide, author)

    # Save output
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
